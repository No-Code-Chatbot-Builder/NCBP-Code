
from PyPDF2 import PdfReader
import io
from docx import Document 
from pptx import Presentation

#tested
async def read_pdf_content(pdf_bytes):
    bytes_buffer = io.BytesIO(pdf_bytes)
    pdf_reader = PdfReader(bytes_buffer)
    text = ''
    for page in pdf_reader.pages:
        text += page.extract_text() if page.extract_text() else ''
    return text

async def read_docx_content(docx_bytes):
    doc = Document(io.BytesIO(docx_bytes))
    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    return text

async def read_pptx_content(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text
            

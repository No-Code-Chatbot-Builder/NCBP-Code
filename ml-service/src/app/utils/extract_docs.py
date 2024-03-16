import fitz
import io
from docx import Document as DocxDocument 
from pptx import Presentation

async def read_pdf_content(pdf_bytes):
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        text = ""
        for page in doc:
            text += page.get_text()
    return text

async def read_docx_content(docx_bytes):
    doc = DocxDocument(io.BytesIO(docx_bytes))
    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    return text

async def read_pptx_content(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text
            
